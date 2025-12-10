import io
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
import base64

def extract_text_from_pdf(file_bytes: bytes, max_pages: int = 20) -> tuple[str, int]:
    """Extract text from PDF bytes."""
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        num_pages = len(doc)
        pages_to_read = min(num_pages, max_pages)
        
        full_text = ""
        for page_num in range(pages_to_read):
            page = doc[page_num]
            text = page.get_text()
            full_text += f"\n--- Page {page_num + 1} ---\n{text}"
        
        doc.close()
        
        # Limit text length to avoid token limits
        if len(full_text) > 50000:
            full_text = full_text[:50000] + "\n\n[Content truncated...]"
        
        return full_text, num_pages
    except Exception as e:
        raise Exception(f"Failed to extract PDF text: {str(e)}")


def extract_text_from_docx(file_bytes: bytes) -> tuple[str, int]:
    """Extract text from DOCX bytes."""
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))
        
        full_text = ""
        page_count = 1
        
        for para in doc.paragraphs:
            full_text += para.text + "\n"
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    full_text += cell.text + " "
                full_text += "\n"
        
        # Estimate page count (roughly 3000 chars per page)
        page_count = max(1, len(full_text) // 3000)
        
        # Limit text length
        if len(full_text) > 50000:
            full_text = full_text[:50000] + "\n\n[Content truncated...]"
        
        return full_text, page_count
    except Exception as e:
        raise Exception(f"Failed to extract DOCX text: {str(e)}")


def extract_text_from_pptx(file_bytes: bytes) -> tuple[str, int]:
    """Extract text from PPTX bytes."""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        
        full_text = ""
        slide_count = len(prs.slides)
        
        for i, slide in enumerate(prs.slides):
            full_text += f"\n--- Slide {i + 1} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text += shape.text + "\n"
        
        # Limit text length
        if len(full_text) > 50000:
            full_text = full_text[:50000] + "\n\n[Content truncated...]"
        
        return full_text, slide_count
    except Exception as e:
        raise Exception(f"Failed to extract PPTX text: {str(e)}")


def extract_text_from_image(file_bytes: bytes, filename: str) -> tuple[str, int]:
    """
    For images, we return a placeholder indicating it's an image.
    The actual image processing will be done by Gemini's vision capability.
    """
    # Return the image as base64 for Gemini to process
    ext = filename.lower().split('.')[-1]
    mime_types = {
        'jpg': 'image/jpeg',
        'jpeg': 'image/jpeg',
        'png': 'image/png',
        'gif': 'image/gif',
        'webp': 'image/webp',
        'bmp': 'image/bmp'
    }
    mime_type = mime_types.get(ext, 'image/jpeg')
    
    # Return a marker that this is an image and include the base64 data
    base64_data = base64.b64encode(file_bytes).decode('utf-8')
    return f"[IMAGE_DATA:{mime_type}:{base64_data}]", 1


def extract_text_from_file(file_bytes: bytes, filename: str) -> tuple[str, int]:
    """Extract text from a file based on its extension."""
    ext = filename.lower().split('.')[-1]
    
    if ext == 'pdf':
        return extract_text_from_pdf(file_bytes)
    elif ext == 'docx':
        return extract_text_from_docx(file_bytes)
    elif ext in ['doc']:
        # Old .doc format - not directly supported, return error message
        return "[This is an older .doc format. Please convert to .docx for better support.]", 1
    elif ext == 'pptx':
        return extract_text_from_pptx(file_bytes)
    elif ext in ['ppt']:
        # Old .ppt format - not directly supported
        return "[This is an older .ppt format. Please convert to .pptx for better support.]", 1
    elif ext in ['jpg', 'jpeg', 'png', 'gif', 'webp', 'bmp']:
        return extract_text_from_image(file_bytes, filename)
    else:
        raise Exception(f"Unsupported file type: {ext}")


def decode_base64_file(base64_data: str) -> bytes:
    """Decode base64 string to bytes, handling data URL prefix."""
    if ',' in base64_data:
        # Remove data URL prefix (e.g., "data:application/pdf;base64,")
        base64_data = base64_data.split(',')[1]
    
    return base64.b64decode(base64_data)
